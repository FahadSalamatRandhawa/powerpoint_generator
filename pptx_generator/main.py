from pptx import Presentation
from pptx.util import Inches
import comtypes.client

def create_ppt(slides_data, ppt_filename):
    prs = Presentation()

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Using a slide layout with title and content
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slide_data['title']
        content.text = slide_data['content']

    prs.save(ppt_filename)

def convert_ppt_to_pdf(ppt_filename, pdf_filename):
    import pypandoc
    pypandoc.download_pandoc()
    output = pypandoc.convert_file(ppt_filename, 'pdf', outputfile=pdf_filename)
    assert output == ""
    st.success(f"PowerPoint presentation converted to PDF and saved as {pdf_filename}")
    # Create a COM object for PowerPoint
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    # Open the PowerPoint presentation
    ppt = powerpoint.Presentations.Open(ppt_filename)

    # Save the presentation as PDF
    ppt.SaveAs(pdf_filename, FileFormat=32)  # 32 is the format code for PDF
    ppt.Close()
    powerpoint.Quit()
    print(f"PowerPoint presentation converted to PDF and saved as {pdf_filename}")

# File paths
ppt_filename = "output_presentation.pptx"
pdf_filename = "output_presentation.pdf"

import fitz
import streamlit as st

def show_pdf(file_path):
    # Display the PDF in Streamlit
    with fitz.open(file_path) as doc:
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap()
            st.image(pix.tobytes(), caption=f'Page {page_num + 1}')